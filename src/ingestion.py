import fitz
import re
import csv
import io
from pathlib import Path
from typing import Generator
from PIL import Image
from langchain_text_splitters import RecursiveCharacterTextSplitter

# ── OCR helpers ──────────────────────────────────────────────────────────────

OCR_THRESHOLD = 50
OCR_CONF_THRESHOLD = 70.0
MIN_TEXT_LENGTH = 50

def _ocr_page(page) -> tuple:
    """Render a PDF page and return (text, avg_confidence) with word-level filtering."""
    import pytesseract
    mat = fitz.Matrix(2, 2)
    pix = page.get_pixmap(matrix=mat)
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

    data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)

    valid_words = []
    valid_confs = []
    low_conf_count = 0

    for i in range(len(data['text'])):
        conf = int(data['conf'][i])
        word = data['text'][i].strip()
        if conf == -1 or not word:
            continue
        if conf >= OCR_CONF_THRESHOLD:
            valid_words.append(word)
            valid_confs.append(conf)
        else:
            low_conf_count += 1

    text = " ".join(valid_words)
    avg_conf = sum(valid_confs) / len(valid_confs) if valid_confs else 0.0
    total_words = len(valid_words) + low_conf_count
    drop_rate = low_conf_count / total_words if total_words > 0 else 0.0

    return text, round(avg_conf, 1), round(drop_rate * 100, 1)

def _ocr_image_bytes(blob: bytes) -> str:
    """Run OCR on raw image bytes (for embedded images in DOCX/PPTX)."""
    import pytesseract
    img = Image.open(io.BytesIO(blob))
    return pytesseract.image_to_string(img)

# ── Format handlers ──────────────────────────────────────────────────────────

def _extract_pdf(path: str) -> Generator[dict, None, None]:
    doc = fitz.open(path)
    for page_num, page in enumerate(doc):
        text = page.get_text()
        is_image_page = len(text.strip()) < OCR_THRESHOLD
        if is_image_page:
            text, avg_conf, drop_rate = _ocr_page(page)
            if len(text.strip()) < MIN_TEXT_LENGTH:
                continue
            yield {
                "text": text,
                "page": page_num + 1,
                "ocr": True,
                "ocr_confidence": avg_conf,
                "ocr_drop_rate": drop_rate,
                "low_quality": avg_conf < OCR_CONF_THRESHOLD
            }
        else:
            if text.strip():
                yield {
                    "text": text,
                    "page": page_num + 1,
                    "ocr": False,
                    "ocr_confidence": 100.0,
                    "ocr_drop_rate": 0.0,
                    "low_quality": False
                }

def _extract_docx(path: str) -> Generator[dict, None, None]:
    from docx import Document
    doc = Document(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    image_texts = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            ocr_text = _ocr_image_bytes(rel.target_part.blob)
            if ocr_text.strip():
                image_texts.append(ocr_text)
    full_text = text + ("\n" + "\n".join(image_texts) if image_texts else "")
    if full_text.strip():
        yield {"text": full_text, "page": 1, "ocr": bool(image_texts),
               "ocr_confidence": 100.0, "ocr_drop_rate": 0.0, "low_quality": False}

def _extract_txt(path: str) -> Generator[dict, None, None]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    if text.strip():
        yield {"text": text, "page": 1, "ocr": False,
               "ocr_confidence": 100.0, "ocr_drop_rate": 0.0, "low_quality": False}

def _extract_pptx(path: str) -> Generator[dict, None, None]:
    from pptx import Presentation
    prs = Presentation(path)
    for slide_num, slide in enumerate(prs.slides):
        texts = [shape.text for shape in slide.shapes
                 if hasattr(shape, "text") and shape.text.strip()]
        for shape in slide.shapes:
            if shape.shape_type == 13:
                ocr_text = _ocr_image_bytes(shape.image.blob)
                if ocr_text.strip():
                    texts.append(ocr_text)
        text = "\n".join(texts)
        if text.strip():
            yield {"text": text, "page": slide_num + 1, "ocr": False,
                   "ocr_confidence": 100.0, "ocr_drop_rate": 0.0, "low_quality": False}

def _extract_csv(path: str) -> Generator[dict, None, None]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.DictReader(f)
        rows = [", ".join(f"{k}: {v}" for k, v in row.items()) for row in reader]
    text = "\n".join(rows)
    if text.strip():
        yield {"text": text, "page": 1, "ocr": False,
               "ocr_confidence": 100.0, "ocr_drop_rate": 0.0, "low_quality": False}

# ── Router ───────────────────────────────────────────────────────────────────

_HANDLERS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".txt":  _extract_txt,
    ".pptx": _extract_pptx,
    ".csv":  _extract_csv,
}

def extract_pages(file_path: str) -> Generator[dict, None, None]:
    ext = Path(file_path).suffix.lower()
    handler = _HANDLERS.get(ext)
    if handler is None:
        raise ValueError(
            f"Unsupported file type: '{ext}'. Supported: {list(_HANDLERS)}"
        )
    yield from handler(file_path)

# ── Text cleaning & chunking ─────────────────────────────────────────────────

def clean_text(text: str) -> str:
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[\x00-\x08\x0b-\x0c\x0e-\x1f]", "", text)
    return text.strip()

def chunk_document(file_path: str, chunk_size: int = 1000, chunk_overlap: int = 100) -> list:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " "]
    )
    chunks = []
    for page_data in extract_pages(file_path):
        cleaned = clean_text(page_data["text"])
        page_chunks = splitter.split_text(cleaned)
        for i, chunk in enumerate(page_chunks):
            chunks.append({
                "text": chunk,
                "page": page_data["page"],
                "chunk_id": f"page{page_data['page']}_chunk{i}",
                "ocr": page_data.get("ocr", False),
                "ocr_confidence": page_data.get("ocr_confidence", 100.0),
                "ocr_drop_rate": page_data.get("ocr_drop_rate", 0.0),
                "low_quality": page_data.get("low_quality", False)
            })
    return chunks
