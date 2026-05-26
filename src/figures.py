import fitz
import io
import hashlib
from pathlib import Path
from PIL import Image

FIGURES_DIR = Path("figures")
FIGURES_DIR.mkdir(exist_ok=True)

MIN_SIZE = 100  # ignore tiny images (icons, bullets) below this px dimension

# ── Helpers ──────────────────────────────────────────────────────────────────

def _save_image(img: Image.Image, source: str, page: int, index: int) -> dict | None:
    """Save a PIL image to disk and return its metadata."""
    if img.width < MIN_SIZE or img.height < MIN_SIZE:
        return None
    stem = Path(source).stem
    fname = f"{stem}_p{page}_fig{index}.png"
    fpath = FIGURES_DIR / fname
    img.save(fpath)
    return {
        "path": str(fpath),
        "filename": fname,
        "source": source,
        "page": page,
        "width": img.width,
        "height": img.height,
    }

# ── Format extractors ────────────────────────────────────────────────────────

def _figures_pdf(path: str) -> list[dict]:
    figures = []
    doc = fitz.open(path)
    for page_num, page in enumerate(doc):
        for idx, img_ref in enumerate(page.get_images(full=True)):
            xref = img_ref[0]
            base = doc.extract_image(xref)
            img = Image.open(io.BytesIO(base["image"]))
            meta = _save_image(img, path, page_num + 1, idx)
            if meta:
                figures.append(meta)
    return figures

def _figures_docx(path: str) -> list[dict]:
    from docx import Document
    figures = []
    doc = Document(path)
    idx = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            img = Image.open(io.BytesIO(rel.target_part.blob))
            meta = _save_image(img, path, 1, idx)
            if meta:
                figures.append(meta)
            idx += 1
    return figures

def _figures_pptx(path: str) -> list[dict]:
    from pptx import Presentation
    figures = []
    prs = Presentation(path)
    for slide_num, slide in enumerate(prs.slides):
        idx = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                img = Image.open(io.BytesIO(shape.image.blob))
                meta = _save_image(img, path, slide_num + 1, idx)
                if meta:
                    figures.append(meta)
                idx += 1
    return figures

# ── Router ───────────────────────────────────────────────────────────────────

_EXTRACTORS = {
    ".pdf":  _figures_pdf,
    ".docx": _figures_docx,
    ".pptx": _figures_pptx,
}

def extract_figures(file_path: str) -> list[dict]:
    """Extract all figures from a document. Returns list of figure metadata."""
    ext = Path(file_path).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return []  # TXT and CSV have no figures
    return extractor(file_path)
